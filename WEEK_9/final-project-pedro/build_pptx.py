# -*- coding: utf-8 -*-
#!/usr/bin/env python3
"""
VaultGuard — PowerPoint Generator
Produces a .pptx matching the HTML presentation structure.

Usage:
    pip install python-pptx pillow
    python build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT_FILE = os.path.join(os.path.dirname(__file__), "VaultGuard_Presentation.pptx")

# ── Palette ────────────────────────────────────────────────────────────────
GOLD        = RGBColor(0xC9, 0xA2, 0x27)
GOLD_DIM    = RGBColor(0xA0, 0x78, 0x18)
DARK_BG     = RGBColor(0x0D, 0x0D, 0x0D)
CARD_BG     = RGBColor(0x1A, 0x1A, 0x1A)
CARD2_BG    = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SILVER      = RGBColor(0x99, 0x99, 0x99)
MONEY       = RGBColor(0x52, 0xB7, 0x88)
MONEY_BRT   = RGBColor(0x74, 0xC9, 0x9A)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW      = RGBColor(0xF1, 0xC4, 0x0F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Add a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, DARK_BG)
    return slide


def fill_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_label(slide, text, left, top, width, height=Inches(0.3),
              size=Pt(9), color=GOLD, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Segoe UI"
    return txBox


def add_textbox(slide, text, left, top, width, height,
                size=Pt(12), color=WHITE, bold=False,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Segoe UI"
    return txBox


def add_title(slide, title, subtitle=None):
    """Gold top stripe + large title."""
    # Gold stripe
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=GOLD)
    # Title
    add_textbox(slide, title,
                Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.0),
                size=Pt(36), color=GOLD, bold=True)
    # Gold underline
    add_rect(slide, Inches(0.8), Inches(1.55), Inches(2), Inches(0.04), fill_color=GOLD)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
                    size=Pt(13), color=SILVER)


def add_bullet_card(slide, title, bullets, left, top, width, height,
                    title_color=WHITE):
    add_rect(slide, left, top, width, height, fill_color=CARD_BG,
             line_color=RGBColor(0x2A, 0x2A, 0x2A))
    add_label(slide, title, left + Inches(0.15), top + Inches(0.12),
              width - Inches(0.3), size=Pt(12), color=title_color, bold=True)
    y = top + Inches(0.42)
    for bullet in bullets:
        add_label(slide, f"• {bullet}", left + Inches(0.15), y,
                  width - Inches(0.3), size=Pt(10), color=RGBColor(0xCC, 0xCC, 0xCC))
        y += Inches(0.26)


def add_kpi(slide, value, unit, left, top, width=Inches(2.2), height=Inches(1.0),
            value_color=GOLD):
    add_rect(slide, left, top, width, height, fill_color=CARD_BG,
             line_color=RGBColor(0x2A, 0x2A, 0x2A))
    add_textbox(slide, value,
                left, top + Inches(0.05), width, Inches(0.55),
                size=Pt(30), color=value_color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, unit,
                left, top + Inches(0.6), width, Inches(0.3),
                size=Pt(9), color=SILVER, align=PP_ALIGN.CENTER)


def add_gold_top_stripe(slide):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill_color=GOLD)


def slide_label(slide, text):
    add_label(slide, text.upper(),
              Inches(0.8), Inches(0.35), Inches(6),
              size=Pt(9), color=GOLD, bold=False)


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def build_slide_01_title(prs):
    """Slide 1 — Title"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    # Big diagonal stripe watermark (simulated with a thin shape)
    add_rect(slide, Inches(9), 0, Inches(0.5), SLIDE_H,
             fill_color=RGBColor(0x16, 0x14, 0x00))

    add_textbox(slide, "VAULTGUARD",
                Inches(0.8), Inches(1.6), Inches(11), Inches(1.4),
                size=Pt(72), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "AI BILLING CLASSIFIER",
                Inches(0.8), Inches(2.9), Inches(11), Inches(0.7),
                size=Pt(30), color=WHITE, bold=False, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.5), Inches(3.65), Inches(4.3), Inches(0.04), fill_color=GOLD)
    add_textbox(slide, "Automating invoice billing decisions with GPT-4o",
                Inches(0.8), Inches(3.75), Inches(11), Inches(0.4),
                size=Pt(14), color=SILVER, align=PP_ALIGN.CENTER)
    add_textbox(slide, "By Pedro Biain",
                Inches(0.8), Inches(6.6), Inches(11), Inches(0.35),
                size=Pt(11), color=SILVER, align=PP_ALIGN.CENTER)


def build_slide_02_context(prs):
    """Slide 2 — Company Context / Funnel"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Company Context")
    add_title(slide, "VaultGuard at a glance", "The scale of the problem")

    items = [
        ("🏦", "100,000", "Smart safes deployed"),
        ("📞", "40,000",  "Service calls / year"),
        ("🔧", "8,000",   "Technician visits / year"),
        ("📋", "8,000",   "Billing decisions / year"),
    ]
    x = Inches(0.8)
    for icon, num, label in items:
        add_rect(slide, x, Inches(2.0), Inches(2.8), Inches(3.5),
                 fill_color=CARD_BG, line_color=GOLD_DIM)
        add_textbox(slide, icon,  x, Inches(2.2),  Inches(2.8), Inches(0.6),
                    size=Pt(28), align=PP_ALIGN.CENTER)
        add_textbox(slide, num,   x, Inches(2.9),  Inches(2.8), Inches(0.8),
                    size=Pt(36), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x, Inches(3.75), Inches(2.8), Inches(0.5),
                    size=Pt(11), color=SILVER, align=PP_ALIGN.CENTER)
        x += Inches(3.05)

    add_textbox(slide,
                "Each billing decision: should the client be charged for this maintenance call or does VaultGuard absorb the cost?",
                Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.4),
                size=Pt(11), color=SILVER)


def build_slide_03_problem(prs):
    """Slide 3 — Problem"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "The Problem")
    add_title(slide, "667 decisions a month — made by hand")

    # KPIs
    kpis = [
        ("667", "invoices / month",     WHITE),
        ("80h",  "manual work / month", WHITE),
        ("7 min","per invoice",         WHITE),
        ("?",    "consistency rate",    RED),
    ]
    x = Inches(0.8)
    for val, lbl, col in kpis:
        add_kpi(slide, val, lbl, x, Inches(2.0), value_color=col)
        x += Inches(2.4)

    # Two issue cards
    add_bullet_card(slide, "Process today",
        ["Biller reads each PDF invoice manually",
         "Checks OOS category against contract rules",
         "Enters YES/NO + GL code in Oracle",
         "No audit trail — decisions vary by day"],
        Inches(0.8), Inches(3.3), Inches(5.5), Inches(2.5))

    add_bullet_card(slide, "Why it matters",
        ["5% error rate -> $100K/yr in missed billings",
         "80 hrs/month of high-skill manual work",
         "No consistency check — decisions vary by individual",
         "Zero visibility for ops or finance teams"],
        Inches(6.7), Inches(3.3), Inches(5.8), Inches(2.5))


def build_slide_04_solution(prs):
    """Slide 4 — Solution"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "The Solution")
    add_title(slide, "AI-powered billing in under 2 seconds")

    # Flow steps
    steps = [
        ("📧", "Receive",  "Email with PDF invoices arrives"),
        ("👁",  "Extract", "GPT-4o reads image-based PDF"),
        ("🧠", "Classify", "YES/NO + GL code + reason"),
        ("⚖️", "Judge",   "Second LLM audits every decision"),
        ("📊", "Deliver",  "Two CSVs emailed to billing team"),
    ]
    x = Inches(0.8)
    for icon, label, desc in steps:
        add_rect(slide, x, Inches(2.0), Inches(2.1), Inches(1.8),
                 fill_color=CARD_BG, line_color=RGBColor(0x2A,0x2A,0x2A))
        add_textbox(slide, icon,  x, Inches(2.1),  Inches(2.1), Inches(0.5),
                    size=Pt(22), align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x, Inches(2.65), Inches(2.1), Inches(0.4),
                    size=Pt(13), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,  x, Inches(3.1),  Inches(2.1), Inches(0.5),
                    size=Pt(9),  color=SILVER, align=PP_ALIGN.CENTER)
        x += Inches(2.4)

    # Before / After time comparison
    add_rect(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(2.4),
             fill_color=CARD_BG, line_color=RGBColor(0x2A,0x2A,0x2A))
    add_textbox(slide, "Time spent on invoice classification — before vs after",
                Inches(1.0), Inches(4.2), Inches(11), Inches(0.35),
                size=Pt(12), color=WHITE, bold=True)
    add_textbox(slide, "Before:  80 hrs / month   (4 hrs/day × 20 working days)",
                Inches(1.0), Inches(4.65), Inches(5.5), Inches(0.35),
                size=Pt(13), color=RED, bold=True)
    add_textbox(slide, "After:   7 hrs / month    (AI handles 92% — human reviews flagged only)",
                Inches(1.0), Inches(5.1), Inches(10), Inches(0.35),
                size=Pt(13), color=MONEY_BRT, bold=True)
    add_textbox(slide, "91% reduction in manual classification work",
                Inches(1.0), Inches(5.55), Inches(10), Inches(0.35),
                size=Pt(11), color=SILVER)


def build_slide_05_poc(prs):
    """Slide 5 — POC Demo"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Proof of Concept")
    add_title(slide, "What the POC demonstrates")

    add_bullet_card(slide, "What it does",
        ["Reads image-based PDF invoices via GPT-4o vision",
         "Extracts 18 structured fields per invoice",
         "Classifies YES/NO with confidence + key signal",
         "Second LLM independently audits every decision",
         "Delivers two formatted spreadsheets by email",
         "Logs all decisions to LangSmith for tracking",
         "GDPR: PII pseudonymised before every LLM call"],
        Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.8))

    kpis = [
        ("92%",   "accuracy vs human",    GOLD),
        ("0",     "dangerous errors",     MONEY_BRT),
        ("2×",    "LLM calls / invoice",  WHITE),
        ("<2s",   "per classification",   WHITE),
    ]
    x = Inches(6.7)
    y = Inches(2.0)
    for val, lbl, col in kpis:
        add_kpi(slide, val, lbl, x, y, width=Inches(2.9), value_color=col)
        y += Inches(1.15)

    add_textbox(slide, "▶  See live n8n workflow demo in the HTML presentation",
                Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.35),
                size=Pt(11), color=SILVER)


def build_slide_06_roi(prs):
    """Slide 6 — ROI"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Return on Investment")
    add_title(slide, "The financial case")

    # Table header
    cols = [Inches(0.8), Inches(6.5), Inches(9.3)]
    headers = ["Cost Item", "Manual (Today)", "With AI"]
    header_y = Inches(2.0)
    add_rect(slide, Inches(0.8), header_y, Inches(11.5), Inches(0.36),
             fill_color=CARD2_BG)
    for i, (col, h) in enumerate(zip(cols, headers)):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        add_textbox(slide, h, col, header_y + Inches(0.04), Inches(2.5), Inches(0.28),
                    size=Pt(10), color=SILVER, bold=True, align=align)

    rows = [
        ("Invoices / month",                          "667",         "667"),
        ("Manual review / month",                     "80 hrs",      "7 hrs"),
        ("Biller cost / year ($50/hr)",               "$50,000",     "$4,200"),
        ("Missed billings / year (5% error rate)",    "~$100,000",   "~$40,000"),
        ("AI hosting (Railway)",                       "—",           "$600"),
        ("Net annual savings",                         "—",           "$105,200"),
    ]
    row_colors = [WHITE, WHITE, WHITE, WHITE, WHITE, MONEY_BRT]
    y = Inches(2.4)
    for i, (label, manual, ai) in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else DARK_BG
        add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.38), fill_color=bg)
        col_val = row_colors[i]
        add_textbox(slide, label,  Inches(0.85), y+Inches(0.04), Inches(5.5),  Inches(0.3),
                    size=Pt(10), color=WHITE if i<5 else MONEY_BRT, bold=(i==5))
        add_textbox(slide, manual, Inches(6.2),  y+Inches(0.04), Inches(2.8),  Inches(0.3),
                    size=Pt(10), color=RED if i in (2,3) else col_val,
                    bold=(i==5), align=PP_ALIGN.RIGHT)
        add_textbox(slide, ai,     Inches(9.1),  y+Inches(0.04), Inches(3.0),  Inches(0.3),
                    size=Pt(10), color=MONEY_BRT if i>=2 else col_val,
                    bold=(i==5), align=PP_ALIGN.RIGHT)
        y += Inches(0.4)

    # Summary KPIs
    kpis = [
        ("$105K",  "saved in year 1",      MONEY_BRT),
        ("$315K",  "saved over 3 years",   MONEY_BRT),
        ("~0",     "months to break even", WHITE),
        ("91%",    "less manual work",     WHITE),
    ]
    x = Inches(0.8)
    for val, lbl, col in kpis:
        add_kpi(slide, val, lbl, x, Inches(5.6), width=Inches(2.7), value_color=col)
        x += Inches(2.9)


def build_slide_07_risks(prs):
    """Slide 7 — Risks"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Risk Assessment")
    add_title(slide, "Risks and mitigations")

    risks = [
        ("⚠️ AI Misclassification",
         ["8% error rate on current eval set",
          "LLM-as-judge catches 80% of classifier errors",
          "All LOW confidence + DISAGREE verdicts flagged for human review",
          "Monthly accuracy benchmarks via evaluate_historical.py"]),
        ("🔒 Data Privacy (GDPR)",
         ["Customer PII pseudonymised before any LLM call",
          "SHA-256 consistent token — same client always same ID",
          "Only extraction step (PDF->fields) sends raw data to OpenAI",
          "All other LLM calls receive sanitised invoice data only"]),
        ("🔌 Vendor Lock-in",
         ["OpenAI model can be swapped — prompt is model-agnostic",
          "n8n workflows are portable (JSON export)",
          "LangSmith traces exportable to CSV",
          "FastAPI + Railway replaceable with any cloud host"]),
        ("📈 Accuracy Drift",
         ["Billing rules change -> system prompt updated centrally",
          "LangSmith dataset grows with each live invoice processed",
          "Historical eval re-run on demand to detect drift",
          "Judge disagreement rate is the early-warning signal"]),
    ]
    positions = [
        (Inches(0.8),  Inches(2.0)),
        (Inches(6.7),  Inches(2.0)),
        (Inches(0.8),  Inches(4.7)),
        (Inches(6.7),  Inches(4.7)),
    ]
    for (risk_title, bullets), (lx, ly) in zip(risks, positions):
        add_bullet_card(slide, risk_title, bullets, lx, ly, Inches(5.6), Inches(2.4))


def build_slide_08_compliance(prs):
    """Slide 8 — Compliance"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Compliance & Legal")
    add_title(slide, "Built with compliance in mind")

    items = [
        ("🔐 GDPR Article 25", "Privacy by Design",
         "PII pseudonymised with SHA-256 before every LLM call. Customer name, address, and contact details never leave the system boundary. Only the extraction step sends raw PDF images to OpenAI — all subsequent steps operate on sanitised data."),
        ("⚖️ Human Oversight", "No fully automated decisions",
         "Every LOW-confidence or judge-DISAGREE invoice is escalated to a human reviewer. The AI is advisory — final decisions remain with the biller. This satisfies GDPR Article 22 (right not to be subject to fully automated decisions)."),
        ("📋 Audit Trail", "Full traceability",
         "Every classification logged to LangSmith with: invoice ID, AI decision, confidence, judge verdict, and human decision (when known). Queryable dataset maintained for compliance review or dispute resolution."),
        ("🏢 Contract Alignment", "White Glove Service Contract",
         "Billing rules are derived directly from VaultGuard's service contract. System prompt reviewed and approved by billing team before deployment. Any contract change triggers a prompt update and re-evaluation."),
    ]
    x = Inches(0.8)
    y = Inches(2.0)
    for i, (icon_title, subtitle, body) in enumerate(items):
        lx = Inches(0.8) if i % 2 == 0 else Inches(6.7)
        ly = Inches(2.0) if i < 2 else Inches(4.6)
        add_rect(slide, lx, ly, Inches(5.6), Inches(2.2),
                 fill_color=CARD_BG, line_color=RGBColor(0x2A,0x2A,0x2A))
        add_textbox(slide, icon_title,
                    lx+Inches(0.15), ly+Inches(0.12), Inches(5.2), Inches(0.35),
                    size=Pt(12), color=GOLD, bold=True)
        add_textbox(slide, subtitle,
                    lx+Inches(0.15), ly+Inches(0.5), Inches(5.2), Inches(0.3),
                    size=Pt(10), color=SILVER)
        add_textbox(slide, body,
                    lx+Inches(0.15), ly+Inches(0.82), Inches(5.3), Inches(1.2),
                    size=Pt(9), color=RGBColor(0xCC,0xCC,0xCC), word_wrap=True)


def build_slide_09_plan(prs):
    """Slide 9 — Strategic Plan"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Strategic Plan")
    add_title(slide, "Three-phase deployment")

    phases = [
        ("Phase 1 — Parallel Run", "Months 1–3",
         ["AI runs alongside current biller",
          "All AI decisions logged vs human decisions",
          "Accuracy target: >90% before go-live",
          "Biller works as normal — AI in shadow mode"]),
        ("Phase 2 — Supervised Live", "Months 4–6",
         ["AI takes all HIGH-confidence decisions",
          "Human reviews MEDIUM / LOW / flagged only",
          "~80% reduction in manual review volume",
          "LangSmith dashboard visible to biller + ops"]),
        ("Phase 3 — Full Automation", "Month 7+",
         ["Human reviews flagged cases only (<10%)",
          "Monthly accuracy benchmark vs ground truth",
          "Quarterly billing rule review + prompt update",
          "Dashboard visible to billing, ops, and finance"]),
    ]
    x = Inches(0.8)
    for phase_title, timeline, bullets in phases:
        add_rect(slide, x, Inches(2.0), Inches(3.8), Inches(4.2),
                 fill_color=CARD_BG, line_color=GOLD_DIM)
        # Top accent bar
        add_rect(slide, x, Inches(2.0), Inches(3.8), Inches(0.06), fill_color=GOLD)
        add_textbox(slide, phase_title,
                    x+Inches(0.15), Inches(2.15), Inches(3.5), Inches(0.4),
                    size=Pt(13), color=GOLD, bold=True)
        add_textbox(slide, timeline,
                    x+Inches(0.15), Inches(2.58), Inches(3.5), Inches(0.3),
                    size=Pt(10), color=SILVER)
        y = Inches(3.0)
        for b in bullets:
            add_label(slide, f"• {b}",
                      x+Inches(0.15), y, Inches(3.5),
                      size=Pt(10), color=RGBColor(0xCC,0xCC,0xCC))
            y += Inches(0.28)
        x += Inches(4.1)


def build_slide_10_mvp(prs):
    """Slide 10 — MVP / What was built"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "MVP — What was built")
    add_title(slide, "Beyond the POC")

    add_bullet_card(slide, "POC -> MVP upgrades",
        ["LLM-as-judge: second AI audits every decision",
         "judge_verdict + judge_reason in API response",
         "GDPR: PII pseudonymised at model boundary",
         "Historical accuracy eval against 89 real invoices",
         "Markdown fence bug fixed (was causing 68% silent failures)",
         "Refined billing rules: coin tube jams + unplugged modems"],
        Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.8))

    add_bullet_card(slide, "Evaluation infrastructure",
        ["evaluate_historical.py — benchmarks AI vs human on demand",
         "3 LangSmith metrics: decision_correct, confidence_appropriate, judge_agrees",
         "No human labelling needed — LLM evaluates LLM"],
        Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.8))

    kpis = [
        ("92%",    "accuracy on 50 historical invoices",            GOLD),
        ("0",      "dangerous errors (HIGH confidence + wrong)",     MONEY_BRT),
        ("2×",     "LLM calls — classifier + independent judge",     WHITE),
        ("$0.004", "cost per invoice classified",                    GOLD),
    ]
    y = Inches(2.0)
    for val, lbl, col in kpis:
        add_kpi(slide, val, lbl, Inches(6.7), y, width=Inches(5.9), value_color=col)
        y += Inches(1.15)


def build_slide_11_roadmap(prs):
    """Slide 11 — Future Product Roadmap"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "What comes next")
    add_title(slide, "Future Product Roadmap")

    cards = [
        ("📱", "Technician Mobile App",
         "Scan safe QR on arrival — see service history, warranty status, correct OOS category. Reduces misclassification at source.",
         "NEW REVENUE"),
        ("📊", "Operations Dashboard",
         "Unified control centre: invoice pipeline, AI accuracy, flagged cases, warranty alerts, fraud hotspots, predictive maintenance queue, SaaS health.",
         "VISIBILITY"),
        ("🔔", "Warranty & Replacement Alerts",
         "Proactive alerts when a safe approaches end-of-warranty. Auto-suspend guarantees before expiry — prevents incorrect NO decisions.",
         "RISK REDUCTION"),
        ("🕵️", "Fraud & Abuse Detection",
         "Pattern analysis across invoices to flag locations with abnormal damage rates — vandalism hotspots, repeated liquid claims, bad OOS logging.",
         "FRAUD PREVENTION"),
        ("🤖", "Predictive Maintenance",
         "Historical invoice data predicts which safes will fail next — by model, age, location. Schedule preventive visits before breakdown.",
         "COST AVOIDANCE"),
        ("🏢", "Multi-Client SaaS",
         "White-label the billing classifier for other cash management operators — ATM networks, parking meters, vending machine fleets.",
         "COMMERCIALISATION"),
    ]
    positions = [
        (Inches(0.8),  Inches(2.0)),
        (Inches(4.7),  Inches(2.0)),
        (Inches(8.6),  Inches(2.0)),
        (Inches(0.8),  Inches(4.7)),
        (Inches(4.7),  Inches(4.7)),
        (Inches(8.6),  Inches(4.7)),
    ]
    for (icon, title, body, badge), (lx, ly) in zip(cards, positions):
        add_rect(slide, lx, ly, Inches(3.6), Inches(2.4),
                 fill_color=CARD_BG, line_color=RGBColor(0x2A,0x2A,0x2A))
        add_rect(slide, lx, ly, Inches(3.6), Inches(0.04), fill_color=GOLD)
        add_textbox(slide, icon,
                    lx+Inches(0.12), ly+Inches(0.1), Inches(0.5), Inches(0.4),
                    size=Pt(18))
        add_textbox(slide, title,
                    lx+Inches(0.12), ly+Inches(0.52), Inches(3.3), Inches(0.35),
                    size=Pt(11), color=WHITE, bold=True)
        add_textbox(slide, body,
                    lx+Inches(0.12), ly+Inches(0.9), Inches(3.3), Inches(1.1),
                    size=Pt(9), color=SILVER, word_wrap=True)
        add_textbox(slide, badge,
                    lx+Inches(0.12), ly+Inches(2.1), Inches(3.3), Inches(0.25),
                    size=Pt(8), color=GOLD)


def build_slide_12_future(prs):
    """Slide 12 — Conclusion"""
    slide = blank_slide(prs)
    add_gold_top_stripe(slide)
    slide_label(slide, "Conclusion")
    add_title(slide, "VaultGuard AI — the case in three numbers")

    kpis = [
        ("92%",   "accuracy vs experienced human biller",     GOLD),
        ("$105K", "net savings in year 1",                    MONEY_BRT),
        ("91%",   "reduction in manual classification work",  WHITE),
    ]
    x = Inches(0.8)
    for val, lbl, col in kpis:
        add_kpi(slide, val, lbl, x, Inches(2.2), width=Inches(3.8), value_color=col)
        x += Inches(4.1)

    add_bullet_card(slide, "Why this works",
        ["GPT-4o understands the nuance in technician notes that simple rules miss",
         "LLM-as-judge eliminates need for human QC on routine invoices",
         "Full audit trail — every decision is traceable, explainable, and reviewable",
         "Accuracy improves over time as the evaluation dataset grows",
         "Zero infrastructure investment — n8n + Railway + OpenAI API"],
        Inches(0.8), Inches(3.6), Inches(5.5), Inches(2.6))

    add_bullet_card(slide, "What was built in this project",
        ["FastAPI on Railway — production REST endpoint",
          "n8n Cloud — end-to-end invoice automation workflow",
          "GPT-4o classifier + LLM-as-judge (2 LLM calls per invoice)",
          "LangSmith tracing + historical evaluation script",
          "Tested against 50 real 2025 invoices: 92% accuracy, 0 dangerous errors"],
        Inches(6.7), Inches(3.6), Inches(5.8), Inches(2.6))


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    print("Building slides...")
    build_slide_01_title(prs)
    print("  OK Slide 1: Title")
    build_slide_02_context(prs)
    print("  OK Slide 2: Context")
    build_slide_03_problem(prs)
    print("  OK Slide 3: Problem")
    build_slide_04_solution(prs)
    print("  OK Slide 4: Solution")
    build_slide_05_poc(prs)
    print("  OK Slide 5: POC Demo")
    build_slide_06_roi(prs)
    print("  OK Slide 6: ROI")
    build_slide_07_risks(prs)
    print("  OK Slide 7: Risks")
    build_slide_08_compliance(prs)
    print("  OK Slide 8: Compliance")
    build_slide_09_plan(prs)
    print("  OK Slide 9: Strategic Plan")
    build_slide_10_mvp(prs)
    print("  OK Slide 10: MVP")
    build_slide_11_roadmap(prs)
    print("  OK Slide 11: Future Roadmap")
    build_slide_12_future(prs)
    print("  OK Slide 12: Conclusion")

    prs.save(OUT_FILE)
    print(f"\nDONE  Saved: {OUT_FILE}")


if __name__ == "__main__":
    main()
